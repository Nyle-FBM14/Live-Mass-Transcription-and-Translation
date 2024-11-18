# Scrapes the daily readings from the Living with Christ website.
# Creates a PowerPoint presentation with slides for each reading.

import requests
from bs4 import BeautifulSoup
from datetime import date, timedelta
from pptx import Presentation
from pptx.util import Inches

#mass_date = date.today() + timedelta(days=3)
mass_date = date.today().strftime("%Y-%m-%d")

def get_reading(reading, date):
    url = ("https://readings.livingwithchrist.ca/daily-texts/%s/%s" %(reading, date))

    web_page = requests.get(url)
    html = BeautifulSoup(web_page.content, "html.parser")
    verse = html.select('h1')[0].text
    return (verse[0:verse.find(" of")], html.select('p')[0].text)

def create_slide(root, title, subtitle, content):
    slide = root.slides.add_slide(root.slide_layouts[0])

    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    #Creates textbox, textframe, then adds content
    slide.shapes.add_textbox(Inches(1), Inches(3), Inches(1), Inches(1)).text_frame.text = content

def create_daily_readings_slides(date = mass_date):
    root = Presentation()

    #First Reading Slide
    (verse, text) = get_reading("reading", mass_date)
    create_slide(root, "First Reading", verse, text)

    #Psalm Slide
    (verse, text) = get_reading("psalm", mass_date)
    create_slide(root, "Psalm", verse, text)

    #Second Reading Slide
    if(not date.weekday()):
        (verse, text) = get_reading("reading2", mass_date)
        create_slide(root, "Second Reading", verse, text)

    #Gospel Slide
    (verse, text) = get_reading("gospel", mass_date)
    create_slide(root, "Gospel", verse, text)

    root.save("Daily_Reading_Slides.pptx")
    print("Done")
    
create_daily_readings_slides()
